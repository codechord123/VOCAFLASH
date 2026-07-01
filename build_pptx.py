#!/usr/bin/env python3
"""행감바·인사약 연습 프레젠테이션 -> PowerPoint(.pptx) 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- 색상 팔레트 (HTML 슬라이드와 동일 톤) ----
BG      = RGBColor(0x0F, 0x17, 0x2A)  # slate-900
CARD    = RGBColor(0x1E, 0x29, 0x3B)  # slate-800
BLUE    = RGBColor(0x60, 0xA5, 0xFA)
BLUE_D  = RGBColor(0x3B, 0x82, 0xF6)
EMER    = RGBColor(0x34, 0xD3, 0x99)
EMER_D  = RGBColor(0x10, 0xB9, 0x81)
RED     = RGBColor(0xF8, 0x71, 0x71)
YELLOW  = RGBColor(0xFD, 0xE0, 0x47)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xCB, 0xD5, 0xE1)  # slate-300
GREY_D  = RGBColor(0x94, 0xA3, 0xB8)  # slate-400

FONT = "Malgun Gothic"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send background to back
    sp = bg._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_run(r, text, size, color, bold=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = FONT


def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    return p


def card(slide, x, y, w, h, fill=CARD, line=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    if line:
        box.line.color.rgb = line; box.line.width = Pt(1.25)
    else:
        box.line.fill.background()
    box.shadow.inherit = False
    return box


def badge(slide, x, y, ch, color):
    d = Inches(0.75)
    b = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background(); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); set_run(r, ch, 20, WHITE, bold=True)
    return b


def pill(slide, x, y, text, txt_color, bg_color):
    w = Inches(2.6)
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.42))
    p.adjustments[0] = 0.5
    p.fill.solid(); p.fill.fore_color.rgb = bg_color
    p.line.fill.background(); p.shadow.inherit = False
    tf = p.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0; tf.margin_bottom = 0
    pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    r = pr.add_run(); set_run(r, text, 12, txt_color, bold=True)
    return p


# =============================================================
# SLIDE 1 — TITLE
# =============================================================
s = add_slide()
pill(s, Inches(4.9), Inches(1.4), "행감바", BLUE, RGBColor(0x1E,0x3A,0x5F))
pill(s, Inches(5.85), Inches(1.4), "인사약", EMER, RGBColor(0x14,0x4D,0x3C))
tf = textbox(s, Inches(1), Inches(2.2), Inches(11.33), Inches(2.2), MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "대화의 두 가지 무기", 44, WHITE, bold=True)
p2 = para(tf, align=PP_ALIGN.CENTER, space_after=0)
set_run(p2.add_run(), "행감바 & 인사약", 52, BLUE, bold=True)
tf2 = textbox(s, Inches(1.5), Inches(4.6), Inches(10.33), Inches(1.2), MSO_ANCHOR.TOP)
p = para(tf2, first=True, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "감정을 다치지 않게 요청하고, 진심이 전해지도록 사과하는 연습", 22, GREY)

# =============================================================
# SLIDE 2 — OVERVIEW
# =============================================================
s = add_slide()
tf = textbox(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(1))
p = para(tf, first=True, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "오늘 연습할 두 가지 공식", 36, WHITE, bold=True)

card(s, Inches(0.9), Inches(1.9), Inches(5.5), Inches(4.6), RGBColor(0x15,0x22,0x3B), BLUE)
tf = textbox(s, Inches(1.25), Inches(2.15), Inches(4.8), Inches(4.1))
p = para(tf, first=True); set_run(p.add_run(), "행감바", 28, BLUE, bold=True)
p = para(tf, space_after=12); set_run(p.add_run(), "불편함을 비난 없이 전하고 원하는 것을 요청할 때", 15, GREY)
for a, b in [("행","동 — 사실만 말한다"),("감","정 — 내 느낌을 말한다"),("바","람 — 원하는 것을 말한다")]:
    p = para(tf, space_after=8)
    set_run(p.add_run(), a, 19, BLUE, bold=True)
    set_run(p.add_run(), b, 19, WHITE)

card(s, Inches(6.9), Inches(1.9), Inches(5.5), Inches(4.6), RGBColor(0x0F,0x2A,0x22), EMER)
tf = textbox(s, Inches(7.25), Inches(2.15), Inches(4.8), Inches(4.1))
p = para(tf, first=True); set_run(p.add_run(), "인사약", 28, EMER, bold=True)
p = para(tf, space_after=12); set_run(p.add_run(), "잘못했을 때 진심이 전해지게 사과할 때", 15, GREY)
for a, b in [("인","정 — 무엇을 잘못했는지 인정"),("사","과 — 진심으로 사과"),("약","속 — 앞으로 어떻게 할지 약속")]:
    p = para(tf, space_after=8)
    set_run(p.add_run(), a, 19, EMER, bold=True)
    set_run(p.add_run(), b, 19, WHITE)

# =============================================================
# Concept slides (행감바 / 인사약)
# =============================================================
def concept_slide(part_label, part_color, part_bg, title, steps, accent):
    s = add_slide()
    pill(s, Inches(5.35), Inches(0.55), part_label, part_color, part_bg)
    tf = textbox(s, Inches(0.8), Inches(1.15), Inches(11.7), Inches(1))
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    set_run(p.add_run(), title, 36, WHITE, bold=True)
    y = Inches(2.35)
    for ch, head, body in steps:
        card(s, Inches(1.6), y, Inches(10.1), Inches(1.15))
        badge(s, Inches(1.9), y + Inches(0.2), ch, accent)
        tf = textbox(s, Inches(2.95), y + Inches(0.12), Inches(8.5), Inches(0.95), MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True, space_after=2)
        set_run(p.add_run(), head, 18, accent, bold=True)
        p = para(tf, space_after=0)
        set_run(p.add_run(), body, 14, GREY)
        y += Inches(1.35)
    return s

concept_slide("PART 1 · 행감바", BLUE, RGBColor(0x1E,0x3A,0x5F), "행동 → 감정 → 바람", [
    ("행","행동 (사실)","평가·비난 없이 '일어난 일'만.  \"네가 항상~\"(X) → \"어제 3시에~\"(O)"),
    ("감","감정 (나-전달법)","\"너 때문에\"(X) → \"나는 ~하게 느꼈어\"(O)"),
    ("바","바람 (구체적 요청)","막연한 불만(X) → \"앞으로 ~해주면 좋겠어\"(O)"),
], BLUE_D)

# =============================================================
# Example slides (좋은 예 / 나쁜 예)
# =============================================================
def example_slide(title, bad_head, bad_text, good_head, good_lines, accent, footer):
    s = add_slide()
    tf = textbox(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(1))
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    set_run(p.add_run(), title, 36, WHITE, bold=True)
    # bad
    card(s, Inches(0.9), Inches(1.8), Inches(5.5), Inches(3.6), RGBColor(0x2E,0x16,0x18), RED)
    tf = textbox(s, Inches(1.25), Inches(2.05), Inches(4.85), Inches(3.1))
    p = para(tf, first=True); set_run(p.add_run(), "✗ " + bad_head, 16, RED, bold=True)
    p = para(tf, space_after=0); p.space_before = Pt(6)
    set_run(p.add_run(), bad_text, 17, WHITE)
    # good
    card(s, Inches(6.9), Inches(1.8), Inches(5.5), Inches(3.6), RGBColor(0x0F,0x2A,0x22), EMER)
    tf = textbox(s, Inches(7.25), Inches(2.05), Inches(4.85), Inches(3.1))
    p = para(tf, first=True); set_run(p.add_run(), "✓ " + good_head, 16, EMER, bold=True)
    for tag, txt in good_lines:
        p = para(tf, space_after=6); p.space_before = Pt(4)
        set_run(p.add_run(), tag + " ", 15, accent, bold=True)
        set_run(p.add_run(), txt, 15, WHITE)
    # footer
    tf = textbox(s, Inches(1), Inches(5.7), Inches(11.33), Inches(1))
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    set_run(p.add_run(), footer, 15, GREY_D)
    return s

example_slide("행감바 예시",
    "이렇게 말하면", "\"너는 왜 맨날 연락도 없이 늦어? 진짜 나 무시하는 거야?\"",
    "행감바로 바꾸면", [
        ("[행]","오늘 약속보다 30분 늦게 왔는데 연락이 없었어."),
        ("[감]","나는 기다리면서 좀 서운하고 걱정됐어."),
        ("[바]","다음엔 늦으면 미리 문자 한 통만 줄 수 있을까?"),
    ], BLUE,
    "핵심: 상대의 인격이 아니라 상황과 내 감정에 대해 말한다.")

# =============================================================
# Practice (빈칸) slides
# =============================================================
def practice_slide(part_label, part_color, part_bg, situation, steps, accent):
    s = add_slide()
    pill(s, Inches(5.35), Inches(0.5), part_label, part_color, part_bg)
    tf = textbox(s, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.8))
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    set_run(p.add_run(), "직접 채워보세요", 32, WHITE, bold=True)
    tf = textbox(s, Inches(1), Inches(1.9), Inches(11.33), Inches(0.6))
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    set_run(p.add_run(), "상황: ", 16, GREY_D)
    set_run(p.add_run(), situation, 16, WHITE, bold=True)
    card(s, Inches(1.6), Inches(2.7), Inches(10.1), Inches(3.6))
    y = Inches(2.95)
    for ch, label in steps:
        badge(s, Inches(2.0), y, ch, accent)
        tf = textbox(s, Inches(3.0), y, Inches(8.3), Inches(1))
        p = para(tf, first=True, space_after=2)
        set_run(p.add_run(), label, 15, accent, bold=True)
        # dashed line
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), y + Inches(0.62), Inches(8.2), Pt(0))
        ln.line.color.rgb = GREY_D; ln.line.width = Pt(1.25)
        ln.line.dash_style = 2
        ln.fill.background(); ln.shadow.inherit = False
        y += Inches(1.05)
    tf = textbox(s, Inches(1), Inches(6.55), Inches(11.33), Inches(0.6))
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    set_run(p.add_run(), "소리 내어 말해본 뒤 다음 슬라이드에서 예시 답안을 확인하세요.", 13, GREY_D)
    return s

practice_slide("연습 시간 · 행감바", BLUE, RGBColor(0x1E,0x3A,0x5F),
    "팀 프로젝트에서 동료가 자기 파트를 마감까지 안 냈다.", [
        ("행","행동 (사실만)"),
        ("감","감정 (나는 ~게 느꼈어)"),
        ("바","바람 (앞으로 ~해주면)"),
    ], BLUE_D)

# =============================================================
# Answer slides
# =============================================================
def answer_slide(lines, accent, footer):
    s = add_slide()
    tf = textbox(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(1))
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    set_run(p.add_run(), "💡 예시 답안", 32, YELLOW, bold=True)
    card(s, Inches(1.6), Inches(2.0), Inches(10.1), Inches(3.3), RGBColor(0x0F,0x2A,0x22), EMER)
    tf = textbox(s, Inches(2.1), Inches(2.4), Inches(9.1), Inches(2.7), MSO_ANCHOR.MIDDLE)
    for i,(tag, txt) in enumerate(lines):
        p = para(tf, first=(i==0), space_after=14)
        set_run(p.add_run(), tag + " ", 19, accent, bold=True)
        set_run(p.add_run(), txt, 19, WHITE)
    tf = textbox(s, Inches(1), Inches(5.6), Inches(11.33), Inches(1))
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    set_run(p.add_run(), footer, 15, GREY_D)
    return s

answer_slide([
    ("[행]","이번 발표 자료에서 네 파트가 마감일까지 올라오지 않았어."),
    ("[감]","나는 마감을 못 맞출까 봐 많이 불안하고 부담됐어."),
    ("[바]","다음부터는 어려우면 하루 전에라도 미리 알려주면 같이 조정할 수 있을 것 같아."),
], BLUE, "정답은 없습니다. 사실 → 내 감정 → 구체적 요청 순서만 지키면 성공!")

# ===== PART 2 : 인사약 =====
concept_slide("PART 2 · 인사약", EMER, RGBColor(0x14,0x4D,0x3C), "인정 → 사과 → 약속", [
    ("인","인정","무엇을 잘못했는지 구체적으로 인정.  \"미안한데~\"(X) → \"내가 ~한 건 잘못이야\"(O)"),
    ("사","사과","변명·\"하지만\" 없이 진심으로. 상대의 감정을 헤아리며."),
    ("약","약속","앞으로 어떻게 다르게 행동할지 구체적 약속."),
], EMER_D)

example_slide("인사약 예시",
    "이런 사과는 역효과", "\"미안, 근데 너도 그럴 만했잖아. 아무튼 미안하다고.\"",
    "인사약으로 바꾸면", [
        ("[인]","내가 사람들 앞에서 네 실수를 지적한 건 잘못이었어."),
        ("[사]","창피하고 속상하게 만들어서 정말 미안해."),
        ("[약]","앞으로 할 말이 있으면 따로 조용히 말할게."),
    ], EMER,
    "금지어: \"하지만\", \"네가 ~했으니까\" — 변명은 사과를 무효로 만든다.")

practice_slide("연습 시간 · 인사약", EMER, RGBColor(0x14,0x4D,0x3C),
    "바빠서 친구의 중요한 부탁을 깜빡하고 못 들어줬다.", [
        ("인","인정 (내가 ~한 건 잘못이야)"),
        ("사","사과 (진심으로, 변명 없이)"),
        ("약","약속 (앞으로 ~할게)"),
    ], EMER_D)

answer_slide([
    ("[인]","네가 부탁한 걸 내가 깜빡하고 안 챙긴 건 명백히 내 잘못이야."),
    ("[사]","너한테 중요한 일이었는데 실망시켜서 진심으로 미안해."),
    ("[약]","앞으로 부탁받으면 바로 메모하고, 오늘 일은 지금 바로 도와줄게."),
], EMER, "진심 어린 사과는 변명 대신 책임, 말 대신 행동으로 끝맺는다.")

# =============================================================
# ROLEPLAY
# =============================================================
s = add_slide()
pill(s, Inches(5.35), Inches(0.55), "롤플레이", YELLOW, RGBColor(0x4D,0x40,0x10))
tf = textbox(s, Inches(0.8), Inches(1.15), Inches(11.7), Inches(1))
p = para(tf, first=True, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "2인 1조 실전 연습", 34, WHITE, bold=True)
cards = [
    ("🕐","상황 카드 1", BLUE, "룸메이트가 설거지를 계속 안 한다. → 행감바로 말하기"),
    ("📱","상황 카드 2", EMER, "약속을 잊고 상대를 바람맞혔다. → 인사약으로 사과하기"),
    ("💬","상황 카드 3", YELLOW, "말다툼 후 화해. → 인사약 먼저, 이어서 행감바"),
]
x = Inches(0.9)
for icon, head, col, body in cards:
    card(s, x, Inches(2.5), Inches(3.85), Inches(2.9))
    tf = textbox(s, x + Inches(0.3), Inches(2.75), Inches(3.3), Inches(2.4))
    p = para(tf, first=True); set_run(p.add_run(), icon, 30, WHITE)
    p = para(tf, space_after=6); set_run(p.add_run(), head, 17, col, bold=True)
    p = para(tf); set_run(p.add_run(), body, 14, GREY)
    x += Inches(4.05)
tf = textbox(s, Inches(1), Inches(5.75), Inches(11.33), Inches(1))
p = para(tf, first=True, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "한 명은 말하고, 한 명은 듣고 피드백. 각 30초씩 말한 뒤 역할을 바꿔보세요.", 15, GREY_D)

# =============================================================
# SUMMARY
# =============================================================
s = add_slide()
tf = textbox(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(1))
p = para(tf, first=True, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "오늘의 정리", 40, WHITE, bold=True)
card(s, Inches(1.2), Inches(2.0), Inches(5.0), Inches(2.5), RGBColor(0x15,0x22,0x3B), BLUE)
tf = textbox(s, Inches(1.5), Inches(2.3), Inches(4.4), Inches(2.0), MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER); set_run(p.add_run(), "행감바 = 요청", 24, BLUE, bold=True)
p = para(tf, align=PP_ALIGN.CENTER); set_run(p.add_run(), "행동 · 감정 · 바람", 20, WHITE)
p = para(tf, align=PP_ALIGN.CENTER); set_run(p.add_run(), "비난하지 않고 내 마음을 전한다", 13, GREY_D)
card(s, Inches(7.1), Inches(2.0), Inches(5.0), Inches(2.5), RGBColor(0x0F,0x2A,0x22), EMER)
tf = textbox(s, Inches(7.4), Inches(2.3), Inches(4.4), Inches(2.0), MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER); set_run(p.add_run(), "인사약 = 사과", 24, EMER, bold=True)
p = para(tf, align=PP_ALIGN.CENTER); set_run(p.add_run(), "인정 · 사과 · 약속", 20, WHITE)
p = para(tf, align=PP_ALIGN.CENTER); set_run(p.add_run(), "변명하지 않고 관계를 회복한다", 13, GREY_D)
tf = textbox(s, Inches(1), Inches(5.0), Inches(11.33), Inches(1.8), MSO_ANCHOR.TOP)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "좋은 대화는 재능이 아니라 연습입니다.", 20, WHITE, bold=True)
p = para(tf, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "오늘 배운 공식을 이번 주 실제 대화에서 한 번 써보세요!", 18, GREY)
p = para(tf, align=PP_ALIGN.CENTER); p.space_before = Pt(14)
set_run(p.add_run(), "감사합니다 🙌", 16, GREY_D)

prs.save("/home/user/VOCAFLASH/행감바_인사약_연습.pptx")
print("saved:", len(prs.slides.__iter__().__length_hint__() if hasattr(prs.slides,'__length_hint__') else 0) or len(prs.slides._sldIdLst), "slides")
